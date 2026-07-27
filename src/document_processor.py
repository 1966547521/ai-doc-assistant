"""Document processing module for AI Document Assistant with logging support.

This module provides comprehensive document reading capabilities for various
file formats including PDF, DOCX, XLSX, PPTX, TXT, and Markdown.
All operations are logged for monitoring and debugging purposes.
"""

import hashlib
import re
from dataclasses import dataclass
from pathlib import Path
from typing import List, Dict, Any

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pypdf import PdfReader
from src.logger import get_logger
from src.semantic_chunker import SemanticChunker

# Initialize module logger
logger = get_logger(__name__)


@dataclass(frozen=True)
class SourceBlock:
    """A cleaned document block with enough metadata to cite its origin."""

    text: str
    metadata: Dict[str, Any]

# Try to import optional document processing libraries
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from openpyxl import load_workbook
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


def clean_text(text: str) -> str:
    """Clean text by removing non-printable characters and normalizing whitespace.

    Args:
        text: Input text to be cleaned

    Returns:
        Cleaned text with normalized whitespace and removed special characters
    """
    if not text:
        return ""

    # Remove BOM and zero-width characters
    text = text.replace("\ufeff", "")
    text = text.replace("\u200b", "").replace("\u200c", "").replace("\u200d", "")
    text = text.replace("\u200e", "").replace("\u200f", "")

    # Remove control characters except newlines and tabs
    cleaned = []
    for char in text:
        if char.isprintable() or char in "\n\r\t ":
            cleaned.append(char)
        elif ord(char) >= 0x10000:  # Remove surrogate pairs and other special Unicode
            cleaned.append(" ")

    text = "".join(cleaned)

    # Remove emoji and special symbols (keep most Unicode text)
    # Keep: Latin (0000-024F), CJK (4E00-9FFF), CJK Ext A (3400-4DBF),
    # Hiragana/Katakana (3040-30FF), Hangul (AC00-D7AF), CJK Ext B+(20000-2A6DF),
    # common punctuation, newlines, tabs
    text = re.sub(
        r"[^\u0000-\u024F\u3400-\u4DBF\u4E00-\u9FFF\u3040-\u30FF\uAC00-\uD7A3"
        r"\u3000-\u303F\uFF00-\uFFEF\U00020000-\U0002A6DF\uF900-\uFAFF"
        r"\U0002F800-\U0002FA1F\n\r\t ]", "", text
    )

    # Remove sequences of special characters
    text = re.sub(
        r"[^\w\u4E00-\u9FFF\u3000-\u303F\uFF00-\uFFEF\s.,!?;:()\[\]{}<>\"\'`~@#$%^&*+-=_|\\/]+", " ", text
    )

    # Normalize whitespace
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"(\n)\s+", r"\1", text)

    # Remove invalid UTF-8 sequences
    text = text.encode("utf-8", errors="replace").decode("utf-8")

    return text.strip()


class DocumentProcessor:
    """Handles reading and splitting documents for the AI assistant.

    Supports multiple document formats: PDF, DOCX, XLSX, PPTX, TXT, and Markdown.
    Automatically detects file type based on extension and extracts text content.
    """

    # Supported file extensions and their descriptions
    SUPPORTED_FORMATS = {
        ".pdf": {"name": "PDF", "description": "Adobe PDF文档"},
        ".docx": {"name": "DOCX", "description": "Microsoft Word文档"},
        ".xlsx": {"name": "XLSX", "description": "Microsoft Excel表格"},
        ".pptx": {"name": "PPTX", "description": "Microsoft PowerPoint演示文稿"},
        ".txt": {"name": "TXT", "description": "纯文本文件"},
        ".md": {"name": "Markdown", "description": "Markdown格式文件"},
    }

    def __init__(self, *, enable_semantic_chunking: bool = False):
        """Initialize the document processor with text splitter."""
        self.semantic_chunker = SemanticChunker() if enable_semantic_chunking else None
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000, chunk_overlap=200, length_function=len
        )

    def read_pdf(self, file_path: str) -> str:
        """Extract text from a PDF file.

        Args:
            file_path: Path to the PDF file

        Returns:
            Extracted and cleaned text content
        """
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            page_text = page.extract_text() or ""
            text += page_text + "\n\n"
        return clean_text(text)

    def read_txt(self, file_path: str) -> str:
        """Read text from a TXT file.

        Args:
            file_path: Path to the TXT file

        Returns:
            Cleaned text content
        """
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return clean_text(f.read())

    def read_markdown(self, file_path: str) -> str:
        """Read text from a Markdown file.

        Args:
            file_path: Path to the Markdown file

        Returns:
            Cleaned text content (Markdown formatting preserved)
        """
        return self.read_txt(file_path)

    def read_docx(self, file_path: str) -> str:
        """Extract text from a DOCX file.

        Args:
            file_path: Path to the DOCX file

        Returns:
            Extracted and cleaned text content

        Raises:
            ImportError: If python-docx library is not installed
        """
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx library not installed. Please install it with: pip install python-docx")

        doc = DocxDocument(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"

        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = "\t".join(cell.text for cell in row.cells)
                text += row_text + "\n"

        return clean_text(text)

    def read_xlsx(self, file_path: str) -> str:
        """Extract text from an XLSX file.

        Args:
            file_path: Path to the XLSX file

        Returns:
            Extracted and cleaned text content with sheet names and cell values

        Raises:
            ImportError: If openpyxl library is not installed
        """
        if not XLSX_AVAILABLE:
            raise ImportError("openpyxl library not installed. Please install it with: pip install openpyxl")

        wb = load_workbook(file_path, read_only=True)
        text = ""

        for sheet_name in wb.sheetnames:
            text += f"【工作表: {sheet_name}】\n"
            sheet = wb[sheet_name]
            
            # Get header row
            headers = []
            if sheet.max_row > 0:
                headers = [cell.value for cell in sheet[1]]
                text += "\t".join(str(h) if h else "" for h in headers) + "\n"
            
            # Get data rows
            for row in sheet.iter_rows(min_row=2, max_row=sheet.max_row):
                row_text = "\t".join(str(cell.value) if cell.value else "" for cell in row)
                text += row_text + "\n"
            text += "\n"

        return clean_text(text)

    def read_pptx(self, file_path: str) -> str:
        """Extract text from a PPTX file.

        Args:
            file_path: Path to the PPTX file

        Returns:
            Extracted and cleaned text content with slide information

        Raises:
            ImportError: If python-pptx library is not installed
        """
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx library not installed. Please install it with: pip install python-pptx")

        prs = Presentation(file_path)
        text = ""

        for i, slide in enumerate(prs.slides, 1):
            text += f"【幻灯片 {i}】\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
                
                # Handle tables in slides
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = "\t".join(cell.text for cell in row.cells)
                        text += row_text + "\n"
            
            text += "\n"

        return clean_text(text)

    def get_supported_formats(self) -> Dict[str, Dict[str, str]]:
        """Get list of supported file formats with availability status.

        Returns:
            Dictionary of supported formats with name, description, and availability
        """
        formats_info = {}
        for ext, info in self.SUPPORTED_FORMATS.items():
            formats_info[ext] = {
                "name": info["name"],
                "description": info["description"],
                "available": self._is_format_available(ext)
            }
        return formats_info

    def _is_format_available(self, ext: str) -> bool:
        """Check if a file format is available (dependencies installed).

        Args:
            ext: File extension to check

        Returns:
            True if format is available, False otherwise
        """
        availability = {
            ".pdf": True,  # pypdf is required dependency
            ".docx": DOCX_AVAILABLE,
            ".xlsx": XLSX_AVAILABLE,
            ".pptx": PPTX_AVAILABLE,
            ".txt": True,
            ".md": True,
        }
        return availability.get(ext, False)

    def read_document(self, file_path: str) -> str:
        """Read document based on file extension.

        Automatically detects file type and extracts text content.

        Args:
            file_path: Path to the document file

        Returns:
            Extracted and cleaned text content

        Raises:
            ValueError: If file format is not supported
            ImportError: If required library is not installed for the format
        """

        lower_path = file_path.lower()
        
        if lower_path.endswith(".pdf"):
            return self.read_pdf(file_path)
        if lower_path.endswith(".docx"):
            return self.read_docx(file_path)
        if lower_path.endswith(".xlsx"):
            return self.read_xlsx(file_path)
        if lower_path.endswith(".pptx"):
            return self.read_pptx(file_path)
        if lower_path.endswith(".txt"):
            return self.read_txt(file_path)
        if lower_path.endswith(".md"):
            return self.read_markdown(file_path)
        
        logger.error(f"Unsupported file format: {file_path}")
        raise ValueError(f"不支持的文件格式: {file_path}")

    def read_source_blocks(self, file_path: str) -> List[SourceBlock]:
        """Extract cleaned blocks together with their location in the source file."""
        lower_path = file_path.lower()
        if lower_path.endswith(".pdf"):
            reader = PdfReader(file_path)
            return [
                SourceBlock(text=text, metadata={"page": page_number})
                for page_number, page in enumerate(reader.pages, start=1)
                if (text := clean_text(page.extract_text() or ""))
            ]

        if lower_path.endswith(".docx"):
            if not DOCX_AVAILABLE:
                raise ImportError(
                    "python-docx library not installed. Please install it with: python-docx"
                )
            doc = DocxDocument(file_path)
            blocks: List[SourceBlock] = []
            for index, paragraph in enumerate(doc.paragraphs, start=1):
                if text := clean_text(paragraph.text):
                    blocks.append(
                        SourceBlock(
                            text=text,
                            metadata={"block_type": "paragraph", "paragraph": index},
                        )
                    )
            for table_index, table in enumerate(doc.tables, start=1):
                for row_index, row in enumerate(table.rows, start=1):
                    if text := clean_text("\t".join(cell.text for cell in row.cells)):
                        blocks.append(
                            SourceBlock(
                                text=text,
                                metadata={
                                    "block_type": "table_row",
                                    "table": table_index,
                                    "row": row_index,
                                },
                            )
                        )
            return blocks

        text = self.read_document(file_path)
        return [SourceBlock(text=text, metadata={})] if text else []

    def split_text(self, text: str) -> List[Document]:
        """Split text into chunks for vector storage.

        Args:
            text: Input text to be split

        Returns:
            List of Document objects, each containing a chunk of text
        """
        chunks = self.text_splitter.split_text(text)
        return [Document(page_content=chunk) for chunk in chunks]

    def split_source_blocks(
        self,
        blocks: List[SourceBlock],
        *,
        document_id: str,
        source_file: str,
    ) -> List[Document]:
        """Split source-aware blocks without losing their citation metadata."""
        chunks: List[Document] = []
        for block in blocks:
            texts = self.text_splitter.split_text(block.text)
            if self.semantic_chunker and re.search(r"[\u4e00-\u9fff]", block.text):
                try:
                    texts = self.semantic_chunker.split(block.text)
                except RuntimeError as exc:
                    logger.warning("Semantic chunking unavailable, using rules: %s", exc)
            for text in texts:
                metadata = {
                    **block.metadata,
                    "document_id": document_id,
                    "source_file": source_file,
                    "chunk_id": f"{document_id}:{len(chunks)}",
                }
                chunks.append(Document(page_content=text, metadata=metadata))
        return chunks

    def process_document(self, file_path: str) -> Dict[str, Any]:
        """Process a document file and return structured information.

        Args:
            file_path: Path to the document file

        Returns:
            Dictionary containing processed document information:
                - text: Cleaned text content
                - chunks: List of Document chunks
                - format: Detected file format
                - char_count: Character count
                - chunk_count: Number of chunks
        """
        logger.debug("Processing document: %s", file_path)

        try:
            # Detect format
            ext = file_path.lower().split(".")[-1]
            format_name = self.SUPPORTED_FORMATS.get(f".{ext}", {}).get("name", ext.upper())

            # Read and process
            source_blocks = self.read_source_blocks(file_path)
            text = "\n\n".join(block.text for block in source_blocks)
            chunks = self.split_source_blocks(
                source_blocks,
                document_id=hashlib.sha256(text.encode("utf-8")).hexdigest(),
                source_file=Path(file_path).name,
            )

            result = {
                "text": text,
                "chunks": chunks,
                "format": format_name,
                "char_count": len(text),
                "chunk_count": len(chunks),
            }

            return result

        except Exception as e:
            logger.error(f"Error processing document: {str(e)}", exc_info=True)
            raise
