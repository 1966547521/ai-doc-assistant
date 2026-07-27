"""Generate sample documents for manual benchmark testing.

Usage:
    cd tests/sample_docs
    python generate_samples.py

Creates .txt, .md files always; also tries .docx / .xlsx / .pptx if libraries available.
"""
import os

OUT_DIR = os.path.dirname(os.path.abspath(__file__))

SECTIONS = [
    ("项目背景", "本项目旨在构建一个基于 RAG（检索增强生成）的智能文档问答系统。"
     "核心链路支持中文 PDF 与 DOCX 的解析、索引和可追溯问答。"),
    ("技术架构", "系统采用分层模块化架构：文档处理层(DocumentProcessor)保留页码与段落来源；"
     "向量存储层(VectorStoreManager)为每份文档创建独立索引；"
     "问答引擎层(QAEngine)使用同一批检索片段生成回答与引用；"
     "Agent 层(AgentSession)采用 Pydantic 结构化意图路由。"),
    ("核心功能", "1. 中文 PDF/DOCX 自动解析与来源保留\n"
     "2. 规则分块与小型本地语义模型分块\n"
     "3. 向量检索与 Top-K 召回（top_k=5）\n"
     "4. 语义缓存（TTL+LRU 淘汰策略）\n"
     "5. Chat 与 Embedding 使用独立、显式的 OpenAI 兼容 API 配置\n"
     "6. 结构化意图 Agent，支持文档问答/摘要/结构分析/信息提取/翻译/报告生成/文档对比"),
    ("Embedding 配置", "系统通过环境变量选择 Embedding 模型：\n"
     "- EMBEDDING_API_KEY / EMBEDDING_BASE_URL / EMBEDDING_MODEL\n"
     "索引按模型与文档隔离，维度不兼容时保留旧索引并明确报错。"),
    ("Agent 工具集", "Agent 使用 Pydantic 结构化意图调用工具，"
     "支持多步协作（如先总结后翻译），路由失败时回到真实 RAG 问答。"),
    ("测试覆盖", "系统提供可重复的单元测试，以及显式启用的真实 API 集成验收。"
     "真实验收会测试 DOCX 解析、向量索引、RAG 问答和来源引用。"),
]


def write_txt():
    path = os.path.join(OUT_DIR, "sample_zh.txt")
    with open(path, "w", encoding="utf-8") as f:
        for title, body in SECTIONS:
            f.write(f"# {title}\n\n{body}\n\n")
    print(f"  {os.path.basename(path)}  — {os.path.getsize(path)} bytes")


def write_md():
    path = os.path.join(OUT_DIR, "sample_zh.md")
    with open(path, "w", encoding="utf-8") as f:
        f.write("---\ntitle: 智能文档助手项目文档\nauthor: Benchmark\n---\n\n")
        for title, body in SECTIONS:
            f.write(f"## {title}\n\n{body}\n\n")
    print(f"  {os.path.basename(path)}  — {os.path.getsize(path)} bytes")


def write_en_txt():
    path = os.path.join(OUT_DIR, "sample_en.txt")
    content = (
        "This is a sample English document for benchmark testing.\n\n"
        + "The quick brown fox jumps over the lazy dog. " * 500
        + "\n\nEnd of document."
    )
    with open(path, "w", encoding="utf-8") as f:
        f.write(content)
    print(f"  {os.path.basename(path)}  — {os.path.getsize(path)} bytes")


def try_write_docx():
    try:
        from docx import Document
        doc = Document()
        doc.add_heading("智能文档助手 - 样本文档", 0)
        for title, body in SECTIONS:
            doc.add_heading(title, 1)
            doc.add_paragraph(body)
        path = os.path.join(OUT_DIR, "sample_zh.docx")
        doc.save(path)
        print(f"  {os.path.basename(path)}  — {os.path.getsize(path)} bytes")
    except ImportError:
        print("  [skip] python-docx not installed")


def try_write_xlsx():
    try:
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws.title = "性能测试"
        headers = ["模块", "功能", "参数", "状态"]
        ws.append(headers)
        rows = [
            ("DocumentProcessor", "多格式解析", "6种格式", "完成"),
            ("VectorStore", "增量索引", "SHA-256去重", "完成"),
            ("QAEngine", "RAG问答", "语义缓存", "完成"),
            ("AgentSession", "ReAct循环", "文本标记调用", "完成"),
        ]
        for row in rows:
            ws.append(row)
        path = os.path.join(OUT_DIR, "sample_zh.xlsx")
        wb.save(path)
        print(f"  {os.path.basename(path)}  — {os.path.getsize(path)} bytes")
    except ImportError:
        print("  [skip] openpyxl not installed")


def try_write_pptx():
    try:
        from pptx import Presentation
        prs = Presentation()
        for title, body in SECTIONS[:4]:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title
            slide.placeholders[1].text = body[:200]
        path = os.path.join(OUT_DIR, "sample_zh.pptx")
        prs.save(path)
        print(f"  {os.path.basename(path)}  — {os.path.getsize(path)} bytes")
    except ImportError:
        print("  [skip] python-pptx not installed")


if __name__ == "__main__":
    print("Generating sample documents...")
    write_txt()
    write_md()
    write_en_txt()
    try_write_docx()
    try_write_xlsx()
    try_write_pptx()
    print("Done.")
